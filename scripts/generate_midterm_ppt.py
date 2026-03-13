from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("docx/mid-term/工业轴承故障预测系统的实现-中期汇报-周逸进.pptx")

NAVY = RGBColor(15, 36, 74)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(95, 168, 211)
CYAN = RGBColor(14, 165, 233)
ORANGE = RGBColor(249, 115, 22)
GREEN = RGBColor(22, 163, 74)
RED = RGBColor(220, 38, 38)
SLATE = RGBColor(71, 85, 105)
LIGHT_BG = RGBColor(247, 249, 252)
CARD_BG = RGBColor(255, 255, 255)
LIGHT_LINE = RGBColor(226, 232, 240)
TEXT = RGBColor(17, 24, 39)
MUTED = RGBColor(100, 116, 139)


def set_background(slide, color: RGBColor = LIGHT_BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, title: str, subtitle: str | None = None) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.95))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(8.6), Inches(0.38))
    paragraph = title_box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.name = "PingFang SC"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(9.2), Inches(0.24), Inches(3.5), Inches(0.28))
        paragraph = sub_box.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        run = paragraph.add_run()
        run.text = subtitle
        run.font.name = "PingFang SC"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(226, 232, 240)


def add_footer(slide, text: str = "中国科学技术大学 软件学院 | 软件工程课程中期检查") -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.07), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_LINE
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.12), Inches(11.8), Inches(0.2))
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "PingFang SC"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED


def add_bullets(slide, left: float, top: float, width: float, height: float, items: list[str], level_indent: float = 0.0) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(8)
        paragraph.bullet = True
        paragraph.left_margin = Pt(level_indent)
        paragraph.font.name = "PingFang SC"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = TEXT


def add_body_text(slide, left: float, top: float, width: float, height: float, lines: list[str], font_size: int = 20, color: RGBColor = TEXT, bold_first: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(8)
        run = paragraph.add_run()
        run.text = line
        run.font.name = "PingFang SC"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if bold_first and index == 0:
            run.font.bold = True


def add_card(slide, left: float, top: float, width: float, height: float, title: str, lines: list[str], accent: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = LIGHT_LINE

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(0.12), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(left + 0.28), Inches(top + 0.16), Inches(width - 0.4), Inches(0.35))
    paragraph = title_box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.name = "PingFang SC"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = accent

    add_body_text(slide, left + 0.28, top + 0.55, width - 0.38, height - 0.7, lines, font_size=15)


def add_progress_row(slide, left: float, top: float, label: str, percent: int, note: str, color: RGBColor) -> None:
    label_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2.3), Inches(0.26))
    paragraph = label_box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = label
    run.font.name = "PingFang SC"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = TEXT

    rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 2.2), Inches(top + 0.02), Inches(4.0), Inches(0.18))
    rail.fill.solid()
    rail.fill.fore_color.rgb = LIGHT_LINE
    rail.line.fill.background()

    fill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 2.2), Inches(top + 0.02), Inches(4.0 * percent / 100), Inches(0.18))
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()

    pct_box = slide.shapes.add_textbox(Inches(left + 6.35), Inches(top - 0.03), Inches(0.7), Inches(0.25))
    paragraph = pct_box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = f"{percent}%"
    run.font.name = "Aptos"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = color

    note_box = slide.shapes.add_textbox(Inches(left + 7.0), Inches(top - 0.03), Inches(5.0), Inches(0.3))
    paragraph = note_box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = note
    run.font.name = "PingFang SC"
    run.font.size = Pt(13)
    run.font.color.rgb = MUTED


def add_timeline_block(slide, left: float, top: float, width: float, title: str, body: str, color: RGBColor) -> None:
    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.18))
    block.fill.solid()
    block.fill.fore_color.rgb = CARD_BG
    block.line.color.rgb = LIGHT_LINE

    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.18), Inches(top + 0.18), Inches(1.18), Inches(0.34))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()

    badge_text = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.21), Inches(0.95), Inches(0.22))
    paragraph = badge_text.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = title
    run.font.name = "PingFang SC"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    add_body_text(slide, left + 1.55, top + 0.16, width - 1.8, 0.8, [body], font_size=15)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = 5765800
    prs.slide_height = 3244850

    # Slide 1: cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.22))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.9), Inches(0.95), Inches(2.15), Inches(2.15))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    circle2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.9), Inches(1.55), Inches(1.25), Inches(1.25))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = CYAN
    circle2.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.0), Inches(8.5), Inches(1.8))
    frame = title_box.text_frame
    p = frame.paragraphs[0]
    r = p.add_run()
    r.text = "工业轴承故障预测系统的实现"
    r.font.name = "PingFang SC"
    r.font.size = Pt(27)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p = frame.add_paragraph()
    r = p.add_run()
    r.text = "软件工程课程中期检查汇报"
    r.font.name = "PingFang SC"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p = frame.add_paragraph()
    p.space_before = Pt(16)
    r = p.add_run()
    r.text = "从开题方案进入工程落地：数据、训练、测试与文档已经形成闭环"
    r.font.name = "PingFang SC"
    r.font.size = Pt(16)
    r.font.color.rgb = RGBColor(191, 219, 254)

    meta = slide.shapes.add_textbox(Inches(0.78), Inches(4.55), Inches(7.4), Inches(1.35))
    frame = meta.text_frame
    for idx, text in enumerate(
        [
            "指导老师：zjf",
            "汇报人：周逸进",
            "小组成员：zyj、cyj、zdh、zy",
            "中国科学技术大学 软件学院 | 2026 年 3 月 13 日",
        ]
    ):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.name = "PingFang SC"
        r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(226, 232, 240)

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.35), Inches(3.25), Inches(0.52))
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(30, 41, 59)
    tag.line.fill.background()
    tag_text = slide.shapes.add_textbox(Inches(0.94), Inches(6.48), Inches(3.0), Inches(0.2))
    p = tag_text.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "关键词：真实数据 | 模块化框架 | 可复现实验"
    r.font.name = "PingFang SC"
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(255, 255, 255)

    # Slide 2: agenda
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "大纲", "Mid-Term Review")
    add_card(slide, 0.8, 1.35, 3.8, 1.2, "01 项目定位", ["目标回顾", "当前阶段位置"], ORANGE)
    add_card(slide, 4.78, 1.35, 3.8, 1.2, "02 已完成工作", ["工程框架", "数据与模型进展"], BLUE)
    add_card(slide, 8.76, 1.35, 3.8, 1.2, "03 结果与验证", ["真实数据 smoke test", "自动化测试"], GREEN)
    add_card(slide, 0.8, 2.85, 3.8, 1.2, "04 当前问题", ["风险点", "应对措施"], RED)
    add_card(slide, 4.78, 2.85, 3.8, 1.2, "05 后续计划", ["第 7-8 周安排", "答辩准备"], CYAN)
    add_card(slide, 8.76, 2.85, 3.8, 1.2, "06 团队协作", ["分工状态", "中期结论"], SKY)
    add_body_text(slide, 0.95, 4.65, 11.3, 1.6, [
        "本次中期汇报沿用开题稿的结构节奏，但内容已经从“方案设计”切换到“工程落地与阶段验证”。",
        "重点展示：真实数据接入、多模型训练链路、回调与实验记录、文档交付体系、当前剩余风险。",
    ], font_size=19)
    add_footer(slide)

    # Slide 3: position
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "项目目标回顾与当前阶段定位", "What We Planned vs. Where We Are")
    add_card(slide, 0.7, 1.35, 4.0, 2.35, "开题阶段承诺", [
        "统一支持 XJTU-SY 与 PHM2012 / FEMTO",
        "构建可扩展的预测框架，而不是单脚本",
        "支持 RUL、阶段分析、可视化与实验记录",
    ], ORANGE)
    add_card(slide, 4.95, 1.35, 3.55, 2.35, "当前阶段位置", [
        "需求与设计基线已成型",
        "核心数据链路与训练主链路已打通",
        "文档体系和 PDF 交付已完成重组",
    ], BLUE)
    add_card(slide, 8.7, 1.35, 3.95, 2.35, "中期核心判断", [
        "项目已经从“方案论证”进入“结果整合”",
        "目前重点不在于再开新坑，而在于收口、稳态和答辩准备",
    ], GREEN)
    add_progress_row(slide, 0.95, 4.35, "需求与文档基线", 100, "proposal 与 mid-term 文档已落地", BLUE)
    add_progress_row(slide, 0.95, 4.8, "数据接入与标签构造", 100, "两个官方数据集已真实验证", GREEN)
    add_progress_row(slide, 0.95, 5.25, "训练、评估与可视化", 90, "主链路完成，继续优化展示细节", ORANGE)
    add_progress_row(slide, 0.95, 5.7, "最终答辩与收尾", 65, "集中在最后两周完成", CYAN)
    add_footer(slide)

    # Slide 4: completed work
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "开题以来已完成的工作", "Completed Since Proposal")
    add_card(slide, 0.7, 1.35, 4.0, 2.45, "工程框架", [
        "建立 USTC.SSE.BearingPrediction 命名空间",
        "统一数据实体、训练器、预测器、评估器接口",
        "形成主程序、测试、文档、脚本的工程组织",
    ], BLUE)
    add_card(slide, 4.95, 1.35, 3.85, 2.45, "数据与特征", [
        "实现 XJTU-SY / PHM2012 Loader",
        "支持滑动窗口、归一化、RMS、峭度、FFT 特征",
        "支持 RUL、阶段标签、健康指标序列构造",
    ], ORANGE)
    add_card(slide, 9.0, 1.35, 3.65, 2.45, "训练与展示", [
        "接入 CNN、RNN、Transformer、MLP",
        "实现 EarlyStopping、TensorBoard、梯度告警",
        "支持预测曲线、阶段图、混淆矩阵、注意力热图",
    ], GREEN)
    add_body_text(slide, 0.9, 4.3, 11.7, 1.7, [
        "阶段性结论：当前系统不是演示脚本，而是一套具备真实数据加载、任务构造、模型训练、结果记录和文档交付能力的课程工程。",
        "与开题相比，最大的变化是“可运行链路”已经被真实数据和测试体系验证。"
    ], font_size=18)
    add_footer(slide)

    # Slide 5: architecture progress
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "系统架构与模块完成度", "Architecture & Module Status")
    add_body_text(slide, 0.95, 1.2, 5.0, 3.2, [
        "数据接入层  ->  预处理与特征工程层",
        "        ->  标签构造层  ->  模型训练与预测层",
        "        ->  评估、可视化与实验管理层",
    ], font_size=24, color=NAVY, bold_first=True)
    add_progress_row(slide, 0.9, 2.6, "M1 数据接入与管理", 100, "XJTU-SY / PHM2012 已跑通", BLUE)
    add_progress_row(slide, 0.9, 3.05, "M2 预处理与特征工程", 100, "原始序列 + 统计特征双入口", ORANGE)
    add_progress_row(slide, 0.9, 3.5, "M3 退化阶段与标签", 100, "3σ、FPT、RUL、HI 序列", GREEN)
    add_progress_row(slide, 0.9, 3.95, "M4 模型训练与预测", 92, "多模型和回调均已工作", CYAN)
    add_progress_row(slide, 0.9, 4.4, "M5 评估、可视化与实验管理", 95, "图表、日志、导出基本齐全", SKY)
    add_card(slide, 7.15, 1.45, 5.35, 3.75, "当前架构特征", [
        "统一对象：BearingEntity / BearingWindowDataset",
        "统一训练：BaseTrainer + Callback 生命周期",
        "统一实验记录：配置、历史、预测、指标、告警",
        "统一输出：CSV / PKL / JSON / YAML / PNG / PT",
    ], NAVY)
    add_footer(slide)

    # Slide 6: data progress
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "数据工程与真实数据进展", "Data Pipeline Progress")
    add_card(slide, 0.75, 1.35, 4.0, 2.6, "官方数据集现状", [
        "XJTU-SY：已下载并完成实体加载验证",
        "PHM2012 / FEMTO：已下载并保留温度通道",
        "原始压缩包已通过 Git LFS 管理",
    ], ORANGE)
    add_card(slide, 4.95, 1.35, 3.75, 2.6, "已解决的真实问题", [
        "CSV 表头和非数值行兼容",
        "温度文件与振动文件按快照对齐",
        "目录层次不一致时的实体路径解析",
    ], BLUE)
    add_card(slide, 8.9, 1.35, 3.7, 2.6, "当前输入能力", [
        "水平振动 / 垂直振动 / 温度",
        "原始序列建模",
        "统计特征建模",
        "健康指标序列建模",
    ], GREEN)
    add_body_text(slide, 0.95, 4.35, 11.2, 1.45, [
        "数据主线是中期阶段最关键的基础工作。当前的训练与测试不再依赖手工拼目录，而是建立在统一 loader 和标签器之上。",
        "这意味着后续新增模型时，不需要重写数据入口。"
    ], font_size=18)
    add_footer(slide)

    # Slide 7: training & experiments
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "训练链路与实验管理进展", "Training, Callbacks and Tracking")
    add_card(slide, 0.75, 1.3, 3.8, 2.9, "模型与预测方式", [
        "CNN / RNN / Transformer / MLP",
        "端到端预测、单步滚动、多步滚动",
        "Monte Carlo Dropout 不确定性估计",
    ], BLUE)
    add_card(slide, 4.75, 1.3, 3.8, 2.9, "训练监控", [
        "EarlyStopping",
        "TensorBoard",
        "GradientAlert",
        "ExperimentLogger",
    ], CYAN)
    add_card(slide, 8.75, 1.3, 3.8, 2.9, "实验记录项", [
        "模型结构、学习率、正则化、epoch 数",
        "采样策略、预测模式、回调配置",
        "history、metrics、predictions、alerts",
    ], GREEN)
    add_body_text(slide, 0.95, 4.6, 11.2, 1.4, [
        "中期阶段的重点不是继续堆模型，而是让“训练过程可监控、实验结果可追踪、异常情况可告警”。",
        "这一点已经基本完成。"
    ], font_size=18)
    add_footer(slide)

    # Slide 8: results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "当前结果与验证情况", "Current Results and Validation")
    add_card(slide, 0.75, 1.3, 3.45, 2.35, "自动化测试", [
        "当前主分支测试：8 passed",
        "覆盖训练回调、数据 IO、阶段划分、预测模式",
        "重点验证真实数据 smoke test 与导出行为",
    ], GREEN)
    add_card(slide, 4.45, 1.3, 3.75, 2.35, "真实数据验证", [
        "PHM2012：可加载加速度与温度通道",
        "XJTU-SY：可解析带表头 CSV",
        "两个官方数据集均完成最小训练/评估验证",
    ], BLUE)
    add_card(slide, 8.45, 1.3, 4.05, 2.35, "交付材料", [
        "11 份课程文档已完成 Markdown + PDF",
        "proposal / mid-term 目录已重组",
        "README 已补充 PDF 导航与说明",
    ], ORANGE)
    add_progress_row(slide, 0.95, 4.3, "代码主链路稳定性", 90, "回调、评估、导出已联通", GREEN)
    add_progress_row(slide, 0.95, 4.75, "真实数据兼容性", 88, "已覆盖主路径，继续补细节", BLUE)
    add_progress_row(slide, 0.95, 5.2, "课程交付完整度", 92, "文档和 PPT 已逐步补齐", ORANGE)
    add_footer(slide)

    # Slide 9: problems
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "当前问题与应对策略", "Issues and Mitigation")
    add_card(slide, 0.75, 1.35, 3.85, 2.9, "问题 1：真实数据体量大", [
        "下载慢、LFS 上传耗时长",
        "不同镜像入口稳定性不一致",
        "处理策略：保留官方入口、本地目录加载与 LFS 原始包",
    ], RED)
    add_card(slide, 4.8, 1.35, 3.85, 2.9, "问题 2：可选依赖差异", [
        "TensorBoard、PyYAML 等环境存在可选性",
        "处理策略：代码中放宽依赖要求，缺失时优雅降级",
        "测试中同步覆盖降级行为",
    ], ORANGE)
    add_card(slide, 8.85, 1.35, 3.75, 2.9, "问题 3：文档口径一致性", [
        "项目名、分工、时间线、路径曾有偏差",
        "处理策略：按 proposal / mid-term 统一重构，并批量导出 PDF",
    ], BLUE)
    add_body_text(slide, 0.95, 4.75, 11.3, 1.2, [
        "当前风险已经从“做不出来”转移到“如何收口得更稳”。因此后半程要优先保证一致性、可展示性和最终答辩效果。"],
        font_size=18)
    add_footer(slide)

    # Slide 10: next steps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "后续工作计划", "Next Two Weeks")
    add_timeline_block(slide, 0.75, 1.45, 5.7, "第 7 周", "补充真实数据实验结果，完善图表输出和答辩演示稿。", BLUE)
    add_timeline_block(slide, 0.75, 2.9, 5.7, "第 8 周", "完成确认测试、最终文档核对、PDF 排版检查和最终提交。", ORANGE)
    add_card(slide, 6.8, 1.45, 5.75, 2.65, "重点收口事项", [
        "把实验结果从“能跑”提升到“可讲清楚”",
        "核对 4 人分工、时间线、术语和目录路径",
        "准备中英文术语统一的答辩表达",
    ], CYAN)
    add_card(slide, 6.8, 4.35, 5.75, 1.7, "明确不再扩张的边界", [
        "不再新增大范围功能，优先打磨已有成果。",
        "不再让文档、代码、演示三套口径脱节。"
    ], GREEN)
    add_footer(slide)

    # Slide 11: team
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, "团队分工与阶段结论", "Team Status")
    add_card(slide, 0.75, 1.35, 3.0, 2.6, "zyj", [
        "系统架构与训练框架",
        "主程序与文档集成",
        "当前重点：收口主线、把控提交质量",
    ], NAVY)
    add_card(slide, 3.95, 1.35, 2.95, 2.6, "cyj", [
        "数据接入、预处理、特征工程",
        "真实数据兼容性修复",
        "当前重点：实验数据与文档一致性",
    ], ORANGE)
    add_card(slide, 7.1, 1.35, 2.95, 2.6, "zdh", [
        "阶段划分、生存分析、确认测试",
        "规格化需求与测试口径",
        "当前重点：验收与评价体系说明",
    ], BLUE)
    add_card(slide, 10.25, 1.35, 2.35, 2.6, "zy", [
        "可视化、图表输出、展示材料",
        "分支修复与交付展示",
        "当前重点：答辩展示效果",
    ], GREEN)
    add_body_text(slide, 0.95, 4.45, 11.2, 1.35, [
        "中期结论：项目的工程骨架已经完整，真实数据、训练流程、测试与文档均已打通。",
        "最后阶段的核心任务是：稳住质量、补强展示、完成最终答辩交付。"
    ], font_size=19)
    add_footer(slide)

    # Slide 12: thanks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.65), Inches(1.0), Inches(2.45), Inches(2.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    circle2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.65), Inches(1.7), Inches(1.2), Inches(1.2))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = CYAN
    circle2.line.fill.background()

    thanks = slide.shapes.add_textbox(Inches(0.95), Inches(1.4), Inches(7.6), Inches(1.4))
    frame = thanks.text_frame
    p = frame.paragraphs[0]
    r = p.add_run()
    r.text = "中期检查汇报完毕"
    r.font.name = "PingFang SC"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p = frame.add_paragraph()
    r = p.add_run()
    r.text = "谢谢老师与各位同学"
    r.font.name = "PingFang SC"
    r.font.size = Pt(24)
    r.font.color.rgb = RGBColor(191, 219, 254)

    sub = slide.shapes.add_textbox(Inches(0.98), Inches(4.1), Inches(8.4), Inches(1.1))
    frame = sub.text_frame
    for idx, text in enumerate([
        "项目名称：工业轴承故障预测系统的实现",
        "课程：软件工程 | 学院：中国科学技术大学 软件学院",
        "文档与代码已统一归档到 proposal / mid-term 目录",
    ]):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.name = "PingFang SC"
        r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(226, 232, 240)

    return prs


def main() -> None:
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUTPUT_PATH)
    print(f"saved {OUTPUT_PATH}")


if __name__ == "__main__":
    main()

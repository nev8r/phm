from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("docx/final/工业轴承故障预测系统的实现-结题答辩.pptx")

NAVY = RGBColor(15, 36, 74)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(14, 165, 233)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(249, 115, 22)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
SLATE = RGBColor(71, 85, 105)
LIGHT_BG = RGBColor(248, 250, 252)
SOFT_BLUE = RGBColor(239, 246, 255)
SOFT_GREEN = RGBColor(240, 253, 244)
SOFT_ORANGE = RGBColor(255, 247, 237)
CARD_BG = RGBColor(255, 255, 255)
LINE = RGBColor(226, 232, 240)
TEXT = RGBColor(17, 24, 39)
MUTED = RGBColor(100, 116, 139)


def set_background(slide, color: RGBColor = LIGHT_BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 16,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: int | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str = "工业轴承故障预测系统的实现 | 结题答辩") -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.78))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    text_box(slide, 0.52, 0.15, 8.8, 0.36, title, size=24, color=RGBColor(255, 255, 255), bold=True)
    text_box(slide, 8.2, 0.25, 4.55, 0.24, subtitle, size=10, color=RGBColor(226, 232, 240), align=PP_ALIGN.RIGHT)


def add_footer(slide, page: int) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.58), Inches(7.08), Inches(12.1), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    text_box(slide, 0.62, 7.13, 9.8, 0.18, "中国科学技术大学 软件学院 | 软件工程课程结题答辩", size=9, color=MUTED)
    text_box(slide, 11.95, 7.13, 0.6, 0.18, f"{page:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, left: float, top: float, width: float, items: list[str], *, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.7))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.bullet = True
        paragraph.space_after = Pt(8)
        paragraph.font.name = "PingFang SC"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = TEXT


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: str,
    color: RGBColor,
    *,
    title_size: int = 16,
    body_size: int = 12,
    fill: RGBColor = CARD_BG,
) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = LINE
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(0.1), Inches(height))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()
    text_box(slide, left + 0.22, top + 0.13, width - 0.34, 0.28, title, size=title_size, color=color, bold=True)
    if body:
        text_box(slide, left + 0.22, top + 0.48, width - 0.36, height - 0.56, body, size=body_size, color=TEXT)


def add_metric(slide, left: float, top: float, value: str, label: str, color: RGBColor) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.35), Inches(1.05))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = LINE
    text_box(slide, left + 0.15, top + 0.18, 2.05, 0.32, value, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, left + 0.15, top + 0.62, 2.05, 0.22, label, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def add_rows(slide, left: float, top: float, rows: list[tuple[str, str, str]], *, row_height: float = 0.54) -> None:
    for index, (left_text, middle_text, right_text) in enumerate(rows):
        y = top + index * row_height
        fill = SOFT_BLUE if index % 2 == 0 else CARD_BG
        row = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(y), Inches(11.8), Inches(row_height - 0.04))
        row.fill.solid()
        row.fill.fore_color.rgb = fill
        row.line.color.rgb = LINE
        text_box(slide, left + 0.15, y + 0.11, 2.75, 0.2, left_text, size=12, color=TEXT, bold=True)
        text_box(slide, left + 3.1, y + 0.11, 4.1, 0.2, middle_text, size=12, color=TEXT)
        text_box(slide, left + 7.45, y + 0.11, 4.0, 0.2, right_text, size=12, color=SLATE)


def add_slide(prs: Presentation, title: str) -> object:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title)
    add_footer(slide, len(prs.slides))
    return slide


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = 5765800
    prs.slide_height = 3244850

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    text_box(slide, 0.7, 1.0, 9.8, 0.72, "工业轴承故障预测系统的实现", size=38, color=RGBColor(255, 255, 255), bold=True)
    text_box(slide, 0.76, 2.0, 9.7, 0.36, "从真实退化数据到 RUL 预测、论文复现与工程闭环", size=20, color=RGBColor(226, 232, 240))
    text_box(slide, 0.78, 5.55, 10.8, 0.32, "中国科学技术大学 软件学院《软件工程》 | 指导老师：zjf", size=14, color=RGBColor(226, 232, 240))
    text_box(slide, 0.78, 5.96, 10.8, 0.32, "小组成员：zyj、cyj、zdh、zy | 2026-06-14", size=14, color=RGBColor(226, 232, 240))

    slide = add_slide(prs, "答辩主线")
    add_metric(slide, 0.8, 1.25, "2", "真实数据集", BLUE)
    add_metric(slide, 3.4, 1.25, "19", "时频域特征", GREEN)
    add_metric(slide, 6.0, 1.25, "2", "RUL 论文复现", ORANGE)
    add_metric(slide, 8.6, 1.25, "31", "自动化测试通过", PURPLE)
    add_bullets(
        slide,
        1.0,
        3.0,
        11.0,
        [
            "项目不是做故障类型分类，而是面向预测性维护的轴承剩余寿命预测系统。",
            "核心价值是把真实数据、特征工程、训练评估、实验记录和课程文档做成一个可运行工程。",
            "答辩重点：我理解数据的时间语义，也能解释模型输入、指标口径和复现边界。",
        ],
        size=17,
    )

    slide = add_slide(prs, "问题定义：为什么是 RUL")
    add_card(slide, 0.8, 1.25, 3.8, 2.0, "业务语义", "轴承不是只回答“是否失效”，而是要估计距离失效还有多久，服务维修窗口安排。", BLUE, fill=SOFT_BLUE)
    add_card(slide, 4.75, 1.25, 3.8, 2.0, "数据语义", "两个数据集都是 run-to-failure 退化过程，天然适合 RUL、健康指标和生存概率分析。", GREEN, fill=SOFT_GREEN)
    add_card(slide, 8.7, 1.25, 3.8, 2.0, "工程语义", "系统输出训练历史、预测表、指标表和图表，便于复现实验和答辩追溯。", ORANGE, fill=SOFT_ORANGE)
    add_bullets(slide, 1.0, 4.05, 11.0, ["因此项目名称和文档统一为“故障预测/剩余寿命预测”，避免把任务误说成故障诊断分类。", "分类能力保留为阶段划分等基础接口，但结题主线是 RUL。"], size=16)

    slide = add_slide(prs, "数据理解：采样组织和时间语义不同")
    add_rows(
        slide,
        0.75,
        1.15,
        [
            ("XJTU-SY", "15 个轴承，3 个工况，每个工况 5 个轴承", "35Hz12kN / 37.5Hz11kN / 40Hz10kN"),
            ("采样方式", "每 1 min 保存一次快照，每个 CSV 32768 点", "25.6 kHz，单快照覆盖 1.28 s"),
            ("PHM2012", "FEMTO/PRONOSTIA 平台，三种工况", "Learning/Test/Full_Test_Set 目录结构"),
            ("采样方式", "每 10 s 一个加速度快照，每个 acc 文件 2560 点", "25.6 kHz，单快照覆盖 0.1 s"),
            ("通道差异", "两者都有水平/垂直振动；PHM2012 额外有温度文件", "loader 统一为 BearingEntity"),
            ("温度边界", "温度字段已对齐保留，便于后续扩展", "本次复现主用振动/时频域特征"),
            ("标签差异", "XJTU-SY 从快照序列推 RUL；PHM2012 Test_set 可叠加已知 terminal RUL", "系统统一成 seconds"),
        ],
    )
    add_card(slide, 0.9, 5.15, 11.5, 0.9, "关键理解", "快照内部是高频振动，快照之间是离散退化时间轴；训练时不能把相邻窗口随机打散，否则容易数据泄漏。", RED, body_size=14)

    slide = add_slide(prs, "特征理解：从振动信号到退化表征")
    add_card(slide, 0.75, 1.1, 3.65, 2.05, "时域特征", "均值、方差、RMS、峰值、峰峰值、峭度、偏度、形状因子、脉冲因子、裕度因子等。", BLUE, fill=SOFT_BLUE)
    add_card(slide, 4.8, 1.1, 3.65, 2.05, "频域特征", "主频、谱能量、谱熵、谱质心、频率 RMS 等，用来观察能量分布和频率迁移。", GREEN, fill=SOFT_GREEN)
    add_card(slide, 8.85, 1.1, 3.65, 2.05, "序列特征", "每个快照提取 19 维特征，再拼成长度 5 或 10 的特征序列供深度模型学习。", ORANGE, fill=SOFT_ORANGE)
    add_bullets(
        slide,
        1.0,
        3.85,
        11.0,
        [
            "RMS、峰值和谱能量通常随退化增强而上升，是健康指标构造的重要基础。",
            "峭度、脉冲因子对冲击性故障更敏感，适合捕捉早期异常波动。",
            "频域特征补充“能量集中在哪些频率”，避免只看振幅大小。",
        ],
        size=16,
    )

    slide = add_slide(prs, "系统架构：工程闭环而不是算法脚本")
    add_rows(
        slide,
        0.75,
        1.1,
        [
            ("data / dataset", "XJTULoader、PHM2012Loader、BearingEntity", "统一真实数据入口"),
            ("preprocess / feature", "鲁棒裁剪、标准化、滑窗、19 维特征", "统一信号处理"),
            ("labeling", "BearingRulLabeler、FeatureSequenceRulLabeler", "统一 RUL 监督目标"),
            ("models / training", "CNN、MLP、Transformer、CNN-LSTM-AM、XLSTM-Transformer", "统一训练接口"),
            ("evaluation", "RMSE、R2、Huang Score、PHM2012 Score", "统一指标口径"),
            ("examples / docs", "notebook、论文复现文档、结项材料", "统一验收交付"),
        ],
    )
    add_card(slide, 0.9, 5.25, 11.5, 0.78, "工程亮点", "模型不直接读原始文件，loader 不混入训练逻辑，实验输出 config/history/metrics/predictions，便于复现和回溯。", PURPLE, body_size=14)

    slide = add_slide(prs, "关键实现：统一数据抽象")
    add_bullets(
        slide,
        0.95,
        1.15,
        11.2,
        [
            "BearingEntity：保存一个轴承的快照表、采样率、工况 metadata 和通道数组。",
            "BearingWindowDataset：保存模型输入、RUL targets、metadata_frame 和 feature_frame。",
            "读前抽样 max_samples：真实数据很大时先等距选文件，再读取 CSV，保证真实训练可控。",
            "时间语义统一：XJTU-SY 1 min 间隔，PHM2012 10 s 间隔，最终 RUL 统一为 seconds。",
        ],
        size=17,
    )
    add_card(slide, 1.15, 5.2, 10.9, 0.76, "答辩可强调", "系统不只是完成 CSV 读取，而是把不同数据集的采样周期、工况、通道和 RUL 单位统一到了后续训练接口。", RED, body_size=14)

    slide = add_slide(prs, "论文复现一：CNN-LSTM-AM")
    add_card(slide, 0.8, 1.05, 3.55, 1.65, "论文核心", "19 维特征 -> CNN 局部编码 -> LSTM 时序建模 -> attention 聚合 -> RUL 回归。", BLUE, fill=SOFT_BLUE)
    add_card(slide, 4.85, 1.05, 3.55, 1.65, "项目实现", "CNNLSTMAttention 支持 use_attention=True/False，直接形成 CNN-LSTM-AM 与 CNN-LSTM 对照。", GREEN, fill=SOFT_GREEN)
    add_card(slide, 8.9, 1.05, 3.55, 1.65, "指标输出", "RMSE、NormalizedRMSE、SMAPE、HuangRulScore、方向性偏差和相对变化列。", ORANGE, fill=SOFT_ORANGE)
    add_rows(
        slide,
        0.85,
        3.25,
        [
            ("XJTU-SY", "Bearing1_5，8 epoch 小样本真实训练", "CNN-LSTM-AM RMSE 406.30"),
            ("PHM2012", "Bearing3_1，8 epoch 小样本真实训练", "CNN-LSTM-AM RMSE 651.04"),
            ("边界", "验证真实训练、指标体系和 workflow 可复现", "不冒充论文完整数值对齐"),
        ],
        row_height=0.62,
    )

    slide = add_slide(prs, "论文复现二：xLSTM-Transformer")
    add_card(slide, 0.8, 1.05, 3.55, 1.75, "为什么选它", "论文同时覆盖 XJTU-SY 和 PHM2012，公开工况划分、序列长度和 RMSE/R2/Score。", BLUE, fill=SOFT_BLUE)
    add_card(slide, 4.85, 1.05, 3.55, 1.75, "实现方式", "XLSTM-Transformer = 指数门控 memory 分支 + Transformer encoder + RUL head。", GREEN, fill=SOFT_GREEN)
    add_card(slide, 8.9, 1.05, 3.55, 1.75, "对比基线", "Feature-Transformer 和 LSTM-Transformer，用对照实验说明结构差异。", ORANGE, fill=SOFT_ORANGE)
    add_rows(
        slide,
        0.85,
        3.35,
        [
            ("XJTU-SY", "3 个工况：4 个轴承训练、1 个轴承测试", "每工况 3 模型"),
            ("PHM2012", "3 个工况：*_1、*_2 训练，*_3 测试", "共 18 行结果"),
            ("边界", "论文结构 + 项目特征管线适配", "非作者源码逐行复刻"),
        ],
        row_height=0.62,
    )

    slide = add_slide(prs, "评价指标：不混淆 Score 口径")
    add_card(slide, 0.75, 1.05, 3.65, 1.9, "普通误差", "MAE、RMSE、NormalizedRMSE、SMAPE、R2：回答预测偏差有多大。", BLUE, fill=SOFT_BLUE)
    add_card(slide, 4.8, 1.05, 3.65, 1.9, "论文 Score", "HuangRulScore 按 Er_i = 100*(target-prediction)/target 分段指数公式。", GREEN, fill=SOFT_GREEN)
    add_card(slide, 8.85, 1.05, 3.65, 1.9, "解释性指标", "Over/UnderPredictionRate、WithinToleranceRate：解释偏早/偏晚预测倾向。", ORANGE, fill=SOFT_ORANGE)
    add_bullets(
        slide,
        1.0,
        3.75,
        11.0,
        [
            "PHM/NASA 类 score 是挑战赛惩罚函数，不等同于 Huang 论文原版 Score。",
            "答辩展示时优先说明指标口径，再解释数值；避免把不同 score 横向误比。",
            "小样本复现的主要价值是工程闭环和指标可落盘，不是追求论文表格完全相同。",
        ],
        size=16,
    )

    slide = add_slide(prs, "测试验收：用证据说话")
    add_metric(slide, 0.85, 1.05, "31", "pytest 全量通过", GREEN)
    add_metric(slide, 3.45, 1.05, "20", "论文与 notebook focused", BLUE)
    add_metric(slide, 6.05, 1.05, "18", "xLSTM 真实训练结果行", ORANGE)
    add_metric(slide, 8.65, 1.05, "8", "真实训练 epoch", PURPLE)
    add_bullets(
        slide,
        1.0,
        2.9,
        11.0,
        [
            "notebook 测试直接执行 examples/*.ipynb，而非只检查文件存在。",
            "真实训练输出 history.csv、metrics.json、predictions.csv、comparison_metrics.csv。",
            "tmp/、outputs/、真实数据和模型产物不提交；文档只记录命令和指标摘要。",
            "四轮提交保存阶段成果，远端 main 已推送，便于回溯。",
        ],
        size=16,
    )

    slide = add_slide(prs, "结项交付物")
    add_rows(
        slide,
        0.75,
        1.08,
        [
            ("源码", "src/USTC/SSE/BearingPrediction", "数据、模型、训练、评估模块"),
            ("示例", "examples/00-07 notebooks", "全部可执行"),
            ("复现文档", "docs/PAPER_REPRODUCTION.md", "两篇论文、命令、指标、边界"),
            ("课程文档", "proposal / mid-term / final", "开题、中期、结题全套材料"),
            ("最终汇报", "docx/final/*.pdf + 结题答辩.pptx", "可直接归档"),
            ("讲稿", "docx/final/md/21_结题答辩演讲稿.md", "逐页答辩稿"),
        ],
    )
    add_card(slide, 0.9, 5.25, 11.5, 0.78, "一句话总结", "项目完成了“真实数据接入 -> RUL 建模 -> 论文复现 -> 自动化测试 -> 课程交付”的完整工程闭环。", RED, body_size=14)

    slide = add_slide(prs, "Q&A 备答：追问与不足")
    add_rows(
        slide,
        0.75,
        1.05,
        [
            ("为什么不是分类？", "数据是 run-to-failure 退化序列", "主任务是 RUL 和预测性维护"),
            ("为什么与论文数值存在差异？", "小样本、8 epoch、CPU 友好验收", "验证流程，不冒充完整论文数值"),
            ("如何避免数据泄漏？", "按时间或轴承划分，notebook 不随机混相邻窗口", "真实复现按论文工况划分"),
            ("xLSTM 是否完全复刻？", "论文无源码，本项目做结构复现", "已在文档说明复现边界"),
            ("工程价值是什么？", "统一数据抽象、训练接口、实验记录、测试文档", "可扩展而非一次性脚本"),
            ("不足与改进？", "小样本 8 epoch，未做多随机种子统计", "后续做全量训练和更强泛化"),
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    text_box(slide, 0.85, 1.2, 10.8, 0.62, "谢谢老师和同学", size=38, color=RGBColor(255, 255, 255), bold=True)
    text_box(slide, 0.9, 2.25, 10.4, 0.36, "欢迎提问：数据语义、特征工程、模型复现、指标口径、工程交付", size=18, color=RGBColor(226, 232, 240))
    add_card(slide, 0.95, 4.5, 10.8, 0.92, "备答底线", "我们完成的是课程项目级真实训练复现和工程闭环，不把小样本结果包装成论文完整数值复现。", ORANGE, body_size=15, fill=CARD_BG)

    return prs


def main() -> None:
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    presentation = build_presentation()
    presentation.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
